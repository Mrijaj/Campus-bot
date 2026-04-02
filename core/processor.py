import os
import fitz  # PyMuPDF
import base64
import httpx
from docx import Document
from pptx import Presentation
from groq import Groq
from config import Config

class FileProcessor:
    def __init__(self):
        # SSL Bypass for HPE/Corporate Network proxy issues
        # verify=False is critical to bypass corporate firewall SSL inspection
        custom_client = httpx.Client(verify=False)
        self.client = Groq(
            api_key=Config.GROQ_API_KEY,
            http_client=custom_client
        )

    def process(self, file_path):
        """Routes files to appropriate extraction logic based on extension."""
        try:
            if not os.path.exists(file_path):
                return ""
            ext = file_path.rsplit('.', 1)[1].lower()

            if ext == 'pdf':
                return self._extract_pdf_smart(file_path)
            elif ext == 'docx':
                return self._extract_docx(file_path)
            elif ext == 'pptx':
                return self._extract_pptx(file_path)
            elif ext in ['jpg', 'jpeg', 'png']:
                return self._extract_image_vision(file_path)
            return ""
        except Exception as e:
            print(f"Error processing {file_path}: {e}")
            return ""

    def _extract_pdf_smart(self, path):
        """Tries digital extraction; falls back to Vision AI if text is sparse or unreadable."""
        text = ""
        try:
            with fitz.open(path) as doc:
                for page in doc:
                    text += page.get_text()
        except Exception as e:
            print(f"Digital PDF read error: {e}")

        # If text is empty (common with scans or complex university encoding),
        # trigger the high-accuracy Llama 4 Vision AI fallback.
        if len(text.strip()) < 50:
            print(f"Digital extraction failed for {path}. Starting Vision Fallback...")
            return self._extract_pdf_via_vision(path)
        return text

    def _extract_pdf_via_vision(self, path):
        """Converts PDF pages to high-res images for Vision OCR."""
        full_text = []
        try:
            with fitz.open(path) as doc:
                # Process up to 3 pages to maintain speed and manage Groq API limits
                for i in range(min(len(doc), 3)):
                    img_path = f"temp_page_{i}.jpg"
                    try:
                        page = doc[i]
                        # 2x Matrix zoom ensures small fee details are legible for AI
                        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                        pix.save(img_path)

                        page_text = self._extract_image_vision(img_path)
                        if page_text:
                            full_text.append(page_text)
                    finally:
                        if os.path.exists(img_path):
                            os.remove(img_path)
            return "\n".join(full_text)
        except Exception as e:
            print(f"PDF-to-Vision conversion failed: {e}")
            return ""

    def _extract_image_vision(self, path):
        """Uses the supported Llama 4 Scout model for high-accuracy OCR."""
        with open(path, "rb") as image_file:
            encoded_image = base64.b64encode(image_file.read()).decode('utf-8')

        try:
            completion = self.client.chat.completions.create(
                # UPDATED MODEL ID: Replacing decommissioned Llama 3.2 Vision
                model="meta-llama/llama-4-scout-17b-16e-instruct",
                messages=[{
                    "role": "user",
                    "content": [
                        {"type": "text", "text": "Extract all text from this image accurately. Return only the extracted text."},
                        {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{encoded_image}"}}
                    ]
                }]
            )
            return completion.choices[0].message.content
        except Exception as e:
            print(f"Vision API Error: {e}")
            return ""

    def _extract_docx(self, path):
        doc = Document(path)
        return "\n".join([p.text for p in doc.paragraphs])

    def _extract_pptx(self, path):
        prs = Presentation(path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)